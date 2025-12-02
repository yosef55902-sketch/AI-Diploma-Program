"""
Extract topics and content from PowerPoint PPTX files
Extracts slide titles, content, and identifies main topics
"""

import json
from pathlib import Path
from pptx import Presentation
import re

def extract_pptx_content(pptx_path):
    """
    Extract content from a PPTX file
    Returns structured data with slides and topics
    """
    try:
        prs = Presentation(pptx_path)
        
        slides_data = []
        all_text = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_data = {
                'slide_number': i,
                'title': '',
                'content': [],
                'full_text': ''
            }
            
            # Extract text from all shapes in the slide
            slide_text_parts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_text_parts.append(text)
                    
                    # First text box is often the title
                    if not slide_data['title'] and len(text) < 200:
                        slide_data['title'] = text
                    else:
                        slide_data['content'].append(text)
            
            slide_data['full_text'] = '\n'.join(slide_text_parts)
            all_text.append(slide_data['full_text'])
            slides_data.append(slide_data)
        
        # Identify main topics from slide titles and content
        topics = identify_topics(slides_data)
        
        return {
            'file_name': pptx_path.name,
            'total_slides': len(slides_data),
            'slides': slides_data,
            'topics': topics,
            'summary': create_summary(slides_data, topics)
        }
    
    except Exception as e:
        return {
            'file_name': pptx_path.name,
            'error': str(e),
            'total_slides': 0,
            'slides': [],
            'topics': [],
            'summary': f'Error extracting content: {e}'
        }

def identify_topics(slides_data):
    """Identify main topics from slides"""
    topics = []
    
    # Extract topics from slide titles
    for slide in slides_data:
        title = slide['title']
        if title:
            # Clean title
            title = title.strip()
            # Remove common prefixes
            title = re.sub(r'^(Slide \d+[:\s]*|Unit \d+[:\s]*)', '', title, flags=re.IGNORECASE)
            if title and len(title) > 3:
                topics.append({
                    'slide': slide['slide_number'],
                    'topic': title,
                    'type': 'title'
                })
    
    # Extract key concepts from content
    key_terms = [
        'ethical framework', 'bias', 'fairness', 'privacy', 'transparency',
        'accountability', 'governance', 'regulation', 'differential privacy',
        'SHAP', 'LIME', 'explainable AI', 'GDPR', 'algorithmic', 'discrimination'
    ]
    
    found_concepts = []
    for slide in slides_data:
        content_lower = slide['full_text'].lower()
        for term in key_terms:
            if term.lower() in content_lower:
                found_concepts.append({
                    'slide': slide['slide_number'],
                    'concept': term,
                    'context': slide['full_text'][:200] + '...' if len(slide['full_text']) > 200 else slide['full_text']
                })
    
    return {
        'titles': topics,
        'key_concepts': found_concepts
    }

def create_summary(slides_data, topics):
    """Create a summary of the PPTX content"""
    if not slides_data:
        return "No content extracted"
    
    summary_parts = []
    
    # Main topics from titles
    if topics['titles']:
        summary_parts.append("Main Topics:")
        for topic in topics['titles'][:5]:  # First 5 topics
            summary_parts.append(f"  - Slide {topic['slide']}: {topic['topic']}")
    
    # Key concepts found
    if topics['key_concepts']:
        concepts = list(set([c['concept'] for c in topics['key_concepts']]))
        summary_parts.append(f"\nKey Concepts Found: {', '.join(concepts[:10])}")
    
    # Slide count
    summary_parts.append(f"\nTotal Slides: {len(slides_data)}")
    
    return '\n'.join(summary_parts)

def extract_all_pptx_files(directory='.'):
    """Extract content from all PPTX files in directory"""
    base_dir = Path(directory)
    pptx_files = sorted(base_dir.glob('*.pptx'))
    
    results = {}
    
    print(f"Found {len(pptx_files)} PPTX files")
    print("="*80)
    
    for pptx_file in pptx_files:
        print(f"\nExtracting: {pptx_file.name}...")
        result = extract_pptx_content(pptx_file)
        results[pptx_file.name] = result
        
        if 'error' in result:
            print(f"  ❌ Error: {result['error']}")
        else:
            print(f"  ✅ Extracted {result['total_slides']} slides")
            print(f"  Topics: {len(result['topics']['titles'])} titles, {len(result['topics']['key_concepts'])} key concepts")
    
    return results

def save_extraction_results(results, output_file='pptx_extraction_results.json'):
    """Save extraction results to JSON file"""
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(results, f, indent=2, ensure_ascii=False)
    print(f"\n✅ Results saved to {output_file}")

def create_markdown_report(results, output_file='PPTX_TOPICS_EXTRACTED.md'):
    """Create a markdown report from extraction results"""
    lines = [
        "# PPTX Topics Extraction - AIAT 116",
        "## استخراج مواضيع ملفات PPTX - AIAT 116",
        "",
        "This document contains extracted topics and content from all PowerPoint presentation files.",
        "",
        "---",
        ""
    ]
    
    for pptx_name in sorted(results.keys()):
        result = results[pptx_name]
        
        lines.append(f"## {pptx_name}")
        lines.append("")
        
        if 'error' in result:
            lines.append(f"**Error:** {result['error']}")
            lines.append("")
            continue
        
        lines.append(f"**Total Slides:** {result['total_slides']}")
        lines.append("")
        
        # Summary
        lines.append("### Summary")
        lines.append("")
        lines.append("```")
        lines.append(result['summary'])
        lines.append("```")
        lines.append("")
        
        # Main Topics
        if result['topics']['titles']:
            lines.append("### Main Topics (from Slide Titles)")
            lines.append("")
            for topic in result['topics']['titles']:
                lines.append(f"- **Slide {topic['slide']}:** {topic['topic']}")
            lines.append("")
        
        # Key Concepts
        if result['topics']['key_concepts']:
            concepts_dict = {}
            for concept in result['topics']['key_concepts']:
                if concept['concept'] not in concepts_dict:
                    concepts_dict[concept['concept']] = []
                concepts_dict[concept['concept']].append(concept['slide'])
            
            lines.append("### Key Concepts Found")
            lines.append("")
            for concept, slides in sorted(concepts_dict.items()):
                slides_str = ', '.join(map(str, slides))
                lines.append(f"- **{concept}:** Found in slides {slides_str}")
            lines.append("")
        
        # Sample slides (first 3)
        if result['slides']:
            lines.append("### Sample Slides Content")
            lines.append("")
            for slide in result['slides'][:3]:
                lines.append(f"#### Slide {slide['slide_number']}")
                if slide['title']:
                    lines.append(f"**Title:** {slide['title']}")
                if slide['content']:
                    lines.append(f"**Content:** {slide['content'][0][:200]}..." if len(slide['content'][0]) > 200 else f"**Content:** {slide['content'][0]}")
                lines.append("")
        
        lines.append("---")
        lines.append("")
    
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines))
    
    print(f"✅ Markdown report saved to {output_file}")

if __name__ == '__main__':
    print("="*80)
    print("PPTX Content Extraction Tool")
    print("="*80)
    
    # Extract from all PPTX files
    results = extract_all_pptx_files()
    
    # Save results
    save_extraction_results(results)
    create_markdown_report(results)
    
    print("\n" + "="*80)
    print("✅ Extraction complete!")
    print("="*80)

